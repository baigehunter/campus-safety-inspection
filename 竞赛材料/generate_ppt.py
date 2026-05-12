"""
Generate competition PPT for 智巡校安
Design: Deep Navy + Cyan Glow, 16:9 format
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
NAVY_DEEP = RGBColor(0x0A, 0x0F, 0x1E)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
NAVY_LIGHT = RGBColor(0x1A, 0x27, 0x44)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
CYAN_GLOW = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_SOFT = RGBColor(0xE2, 0xE8, 0xF0)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GRAY_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)
GREEN_ACCENT = RGBColor(0x10, 0xB9, 0x81)
ORANGE_ACCENT = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)


def add_bg(slide, color=NAVY_DEEP):
    """Solid dark background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_cyan_bar(slide, left, top, width, height):
    """Cyan accent line"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """Add text box with styling"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_rich_text(tf, text, font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Add paragraph to existing text frame"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = alignment
    p.space_after = Pt(6)
    return p


def add_card(slide, left, top, width, height, color=NAVY_LIGHT):
    """Rounded rectangle card"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, color=CYAN):
    """Small decorative circle"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY_DEEP)

# Decorative elements
add_circle(slide, Inches(-1.5), Inches(-1.5), Inches(4), RGBColor(0x06, 0xB6, 0xD4))
add_circle(slide, Inches(10), Inches(5), Inches(5), RGBColor(0x1A, 0x27, 0x44))
add_cyan_bar(slide, Inches(1.5), Inches(3.2), Inches(0.8), Pt(4))

# Title
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             '智巡校安', font_size=56, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
             '基于国产大模型的校园安全智能巡检信息系统', font_size=28, color=CYAN_GLOW)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.6),
             '2026年教师人工智能应用案例征集活动 — 智能信息系统赛道', font_size=16, color=GRAY)

# Bottom info
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(5), Inches(0.5),
             '开发工具：Python/Django + uni-app/Vue 3 + 豆包(火山方舟) + Claude', font_size=14, color=GRAY_LIGHT)

# ═══════════════════════════════════════════════
# SLIDE 2 — Background: Pain Points
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

# Section title
add_text_box(slide, Inches(1.5), Inches(0.7), Inches(5), Inches(0.7),
             '开发背景 — 现实困境', font_size=32, color=WHITE, bold=True)

# Three pain point cards
card_data = [
    ('📋 记录易丢失', '纸质表格打钩签字，归档后难以回溯\n曾发生消防栓漏水却找不到上次巡检记录的案例'),
    ('🔍 判断靠经验', '巡检员多为保安和后勤人员\n同一隐患不同人员判断结果差异大，误报漏报频发'),
    ('🔄 整改缺闭环', '隐患上报后在微信群沟通，信息碎片化\n上学期三起隐患超过两周无人跟进'),
]

for i, (title, desc) in enumerate(card_data):
    left = Inches(1.5 + i * 3.6)
    top = Inches(2.0)
    add_card(slide, left, top, Inches(3.2), Inches(3.8), NAVY_LIGHT)

    # Number circle
    add_circle(slide, left + Inches(1.2), top + Inches(0.5), Inches(0.7), CYAN)
    add_text_box(slide, left + Inches(1.2), top + Inches(0.55), Inches(0.7), Inches(0.6),
                 str(i + 1), font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.3), top + Inches(1.6), Inches(2.6), Inches(0.6),
                 title, font_size=22, color=WHITE, bold=True)

    tf = add_text_box(slide, left + Inches(0.3), top + Inches(2.3), Inches(2.6), Inches(1.3),
                      '', font_size=14, color=GRAY)
    tf.paragraphs[0].text = desc.split('\n')[0]
    add_rich_text(tf, desc.split('\n')[1] if '\n' in desc else '', font_size=14, color=GRAY_LIGHT)

# Bottom note
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
             '这些问题是中小学及高校普遍面临的管理难题，本质上是"数据采集—智能判定—流程闭环"的信息化需求',
             font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 3 — Background: AI Feasibility
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '开发背景 — 借助AI解决问题的可行性', font_size=32, color=WHITE, bold=True)

# Three pillars
pillars = [
    ('🧠', '国产大模型', '豆包（火山方舟）已具备\n图像理解与文本分类能力\nAPI成本极低，合规性强'),
    ('📱', '成熟技术栈', '微信小程序用户零安装\nDjango后端快速构建\n技术生态完善，资料丰富'),
    ('🤖', 'AI辅助编程', 'Claude可有效辅助代码开发\n帮助非专业人员跨越技能门槛\n"AI赋能终身学习"的实践'),
]

for i, (icon, title, desc) in enumerate(pillars):
    left = Inches(1.5 + i * 3.6)
    top = Inches(2.0)
    add_card(slide, left, top, Inches(3.2), Inches(3.5), NAVY_LIGHT)

    add_text_box(slide, left + Inches(0.3), top + Inches(0.4), Inches(2.6), Inches(0.8),
                 icon, font_size=48, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.4), Inches(2.6), Inches(0.5),
                 title, font_size=22, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_text_box(slide, left + Inches(0.3), top + Inches(2.0), Inches(2.6), Inches(1.4),
                      '', font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    lines = desc.split('\n')
    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        add_rich_text(tf, line, font_size=14, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# Bottom emphasis
add_text_box(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.6),
             '非计算机专业教师借助AI，也有能力独立完成一套三端系统的设计与开发',
             font_size=18, color=CYAN_GLOW, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 4 — Tech Selection
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '设计与开发 — 技术选型', font_size=32, color=WHITE, bold=True)

# Tech stack table as cards
tech_data = [
    ('AI分析引擎', '豆包 doubao-seed-2-0-lite（火山方舟）', '国产合规，图像理解强，成本低'),
    ('AI辅助开发', 'Claude（Anthropic）', '代码生成、调试、前端设计辅助'),
    ('后端框架', 'Python 3.13 + Django 6.0 + DRF', 'ORM成熟，JWT认证体系完善'),
    ('数据库', 'MySQL 8.0', '适合结构化巡检台账数据'),
    ('移动端', 'uni-app（Vue 3 + TypeScript）', '一套代码出微信小程序，零安装'),
    ('管理后台', 'Vue 3 + Element Plus + ECharts', '组件丰富，数据可视化强'),
    ('消息推送', '微信订阅消息', '合规触达，无额外成本'),
    ('部署环境', 'Nginx + NSSM（Windows服务）', '适配学校Windows Server环境'),
]

for i, (layer, tech, reason) in enumerate(tech_data):
    row = i // 2
    col = i % 2
    left = Inches(1.5 + col * 5.4)
    top = Inches(1.9 + row * 1.25)

    add_card(slide, left, top, Inches(5.0), Inches(1.05), NAVY_LIGHT)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(1.6), Inches(0.35),
                 layer, font_size=13, color=CYAN, bold=True)
    add_text_box(slide, left + Inches(1.8), top + Inches(0.1), Inches(3.0), Inches(0.35),
                 tech, font_size=13, color=WHITE)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.55), Inches(4.6), Inches(0.4),
                 reason, font_size=11, color=GRAY)

# ═══════════════════════════════════════════════
# SLIDE 5 — Development: AI Prompt Design
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '开发过程 — AI隐患分析的提示词设计', font_size=32, color=WHITE, bold=True)

# Left: Iteration process
add_card(slide, Inches(1.5), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(1.8), Inches(2.1), Inches(4.5), Inches(0.5),
             '提示词迭代过程（AI辅助生成）', font_size=20, color=CYAN, bold=True)

steps = [
    '① 初始方案过于简单 → AI返回不稳定，有时拒答',
    '② 添加角色设定："你是校园安全巡检专家"',
    '③ 约束JSON输出格式：hazard_type / level / tags / analysis',
    '④ 调整temperature = 0.3，降低输出随机性',
    '⑤ 约7轮迭代，测试50个样本，准确率稳定在92%',
]
for i, step in enumerate(steps):
    add_text_box(slide, Inches(1.8), Inches(2.8 + i * 0.7), Inches(4.5), Inches(0.5),
                 step, font_size=14, color=WHITE_SOFT if i == 4 else GRAY)

# Right: Technical solution
add_card(slide, Inches(7.3), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(7.6), Inches(2.1), Inches(4.5), Inches(0.5),
             '照片传输的技术方案', font_size=20, color=CYAN, bold=True)

add_text_box(slide, Inches(7.6), Inches(2.8), Inches(4.5), Inches(0.3),
             '问题：', font_size=15, color=RED_ACCENT, bold=True)
add_text_box(slide, Inches(7.6), Inches(3.1), Inches(4.5), Inches(0.8),
             '微信小程序图片为本地路径(wxfile://)，\n服务器无法直接访问', font_size=14, color=GRAY)

add_text_box(slide, Inches(7.6), Inches(3.9), Inches(4.5), Inches(0.3),
             '方案：', font_size=15, color=GREEN_ACCENT, bold=True)
add_text_box(slide, Inches(7.6), Inches(4.2), Inches(4.5), Inches(1.2),
             'AI分析时：照片 → base64编码 → 发给豆包\n提交上报时：先上传图片 → 获取URL → 再提交表单\n\n区分策略避免base64过大导致超时',
             font_size=14, color=WHITE_SOFT)

add_text_box(slide, Inches(7.6), Inches(5.6), Inches(4.5), Inches(0.4),
             '（提示词框架 + 照片传输方案由AI辅助生成，人工测试迭代）',
             font_size=12, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 6 — Development: Frontend Polish
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '开发过程 — 前端界面与动画系统的打磨', font_size=32, color=WHITE, bold=True)

# Two columns
add_card(slide, Inches(1.5), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(1.8), Inches(2.1), Inches(4.5), Inches(0.5),
             'AI生成 → 人工调优', font_size=20, color=CYAN, bold=True)

items = [
    '35个CSS自定义属性，构建完整设计系统',
    '8组关键帧动画 + 8级延迟入场机制',
    '真机测试：动画过多导致卡顿 → 精简动画',
    '动画时长 0.45s → 0.38s，减少等待感',
    '移除backdrop-filter（微信小程序不支持）',
    '添加prefers-reduced-motion无障碍适配',
]
for i, item in enumerate(items):
    prefix = '🤖' if i < 2 else '👤'
    add_text_box(slide, Inches(1.8), Inches(2.8 + i * 0.65), Inches(4.8), Inches(0.5),
                 f'{prefix} {item}', font_size=14, color=WHITE_SOFT if i >= 2 else GRAY)

add_card(slide, Inches(7.3), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(7.6), Inches(2.1), Inches(4.5), Inches(0.5),
             '可访问性适配', font_size=20, color=CYAN, bold=True)

add_text_box(slide, Inches(7.6), Inches(2.9), Inches(4.5), Inches(1.2),
             '字体大小调整\n28rpx → 32rpx\n\n多位巡检员反馈户外光线下字太小\n逐页调整所有页面字号',
             font_size=15, color=WHITE_SOFT)

add_text_box(slide, Inches(7.6), Inches(4.3), Inches(4.5), Inches(1.2),
             '对比度优化\n#94a3b8 → #64748b\n\n灰色文字在强光下几乎不可见\n加深配色确保户外可见性',
             font_size=15, color=WHITE_SOFT)

add_text_box(slide, Inches(7.6), Inches(5.6), Inches(4.5), Inches(0.4),
             '（CSS框架与动画系统由AI生成，人工完成真机适配）',
             font_size=12, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 7 — Feature Architecture Overview
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '功能架构 — 三端系统总览', font_size=32, color=WHITE, bold=True)

# Flow diagram
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.5),
             '核心业务流程：发现上报 → 智能判定 → 派单整改 → 验收闭环',
             font_size=18, color=CYAN_GLOW, bold=True, alignment=PP_ALIGN.CENTER)

# Three platform cards
platforms = [
    ('📱 移动端小程序', '巡检员 / 整改人', [
        '登录认证 · 点位巡检 · AI判定',
        '隐患上报 · 整改任务 · 历史记录',
        '消息中心 · 个人中心',
    ]),
    ('🖥️ Web管理后台', '安全管理员', [
        '数据看板(ECharts) · 用户管理',
        '区域管理 · 点位管理 · 巡检频率',
        '隐患审核派单 · 整改验收 · 操作日志',
    ]),
    ('⚙️ 后端API', 'Django REST Framework', [
        '30+ RESTful接口 · JWT双Token认证',
        '自定义权限类 · 角色级访问控制',
        'AI分析服务 · 微信推送 · 文件上传',
    ]),
]

for i, (title, role, items) in enumerate(platforms):
    left = Inches(1.5 + i * 3.6)
    top = Inches(3.0)
    add_card(slide, left, top, Inches(3.2), Inches(3.5), NAVY_LIGHT)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.8), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.6), Inches(2.8), Inches(0.3),
                 role, font_size=12, color=CYAN)

    for j, item in enumerate(items):
        add_text_box(slide, left + Inches(0.3), top + Inches(1.2 + j * 0.7), Inches(2.6), Inches(0.5),
                     f'• {item}', font_size=13, color=WHITE_SOFT)

# ═══════════════════════════════════════════════
# SLIDE 8 — Mini Program Features (Detail)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '功能详情 — 移动端小程序', font_size=32, color=WHITE, bold=True)

# 7 features in 2 rows + 1
features = [
    ('🔐', '账号认证', 'JWT双Token\n自动续期\nIP限速防暴破'),
    ('📍', '巡检点位', '点位搜索\n拍照记录\nAI判定状态'),
    ('⚠️', '隐患上报', '结构化表单\nAI自动分类定级\n结果可人工修正'),
    ('📋', '整改任务', '待处理/待验收/已完成\n拍照提交\n驳回可重提'),
    ('📊', '历史记录', '巡检+隐患双Tab\nAI结论回顾\n照片预览'),
    ('🔔', '消息中心', '分类筛选\n未读标记\n下拉刷新'),
    ('👤', '个人中心', '用户信息\n统计数据\n订阅管理'),
]

for i, (icon, title, desc) in enumerate(features):
    row = i // 4
    col = i % 4
    if i == 6:  # last item centered
        left = Inches(4.8)
        top = Inches(5.5)
    else:
        left = Inches(1.5 + col * 2.7)
        top = Inches(1.9 + row * 2.8)

    add_card(slide, left, top, Inches(2.4), Inches(2.4), NAVY_LIGHT)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.2), Inches(2.1), Inches(0.5),
                 icon, font_size=32, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.75), Inches(2.1), Inches(0.4),
                 title, font_size=18, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_text_box(slide, left + Inches(0.15), top + Inches(1.25), Inches(2.1), Inches(0.9),
                      '', font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
    lines = desc.split('\n')
    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        add_rich_text(tf, line, font_size=12, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 9 — Web Admin Features
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '功能详情 — Web管理后台与后端API', font_size=32, color=WHITE, bold=True)

# Left: Web Admin
add_card(slide, Inches(1.5), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(1.8), Inches(2.1), Inches(4.5), Inches(0.5),
             '🖥️ Web管理后台', font_size=20, color=CYAN, bold=True)

admin_modules = [
    ('数据看板', '巡检/隐患统计图表 → ECharts可视化，支持下钻'),
    ('用户管理', '四种角色 → 超级管理员/安全管理员/巡检员/整改负责人'),
    ('区域点位', '校园区域层级 → 巡检频率(日/周/月) → 点位状态'),
    ('隐患管理', '列表查看 → 人工派单 → 状态跟踪 → 闭环归档'),
    ('整改审核', '结果审核 → 验收通过 → 或驳回重提'),
    ('操作日志', '全操作留痕 → 操作人/时间/IP → 审计追溯'),
]
for i, (name, desc) in enumerate(admin_modules):
    add_text_box(slide, Inches(1.8), Inches(2.8 + i * 0.62), Inches(4.8), Inches(0.55),
                 f'• {name}：{desc}', font_size=13, color=WHITE_SOFT)

# Right: Backend API
add_card(slide, Inches(7.3), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(7.6), Inches(2.1), Inches(4.5), Inches(0.5),
             '⚙️ 后端API服务', font_size=20, color=CYAN, bold=True)

add_text_box(slide, Inches(7.6), Inches(2.8), Inches(4.5), Inches(0.8),
             '30+ RESTful接口，覆盖全部业务', font_size=14, color=WHITE, bold=True)

api_categories = [
    '认证模块：登录/登出/Token刷新（JWT + IP限速）',
    '巡检模块：点位CRUD / 记录提交 / AI判定',
    '隐患模块：上报 / AI分析 / 状态流转',
    '整改模块：任务分配 / 提交审核 / 验收闭环',
    '通知模块：消息推送 / 已读标记 / 未读计数',
    '统计模块：仪表盘数据 / 图表数据 / 报表导出',
]
for i, api in enumerate(api_categories):
    add_text_box(slide, Inches(7.6), Inches(3.6 + i * 0.55), Inches(4.8), Inches(0.5),
                 f'• {api}', font_size=13, color=GRAY)

add_text_box(slide, Inches(7.6), Inches(6.1), Inches(4.5), Inches(0.4),
             '权限控制：自定义权限类 + 角色级数据隔离',
             font_size=12, color=CYAN, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 10 — Pilot Deployment
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '应用过程 — 试点部署情况', font_size=32, color=WHITE, bold=True)

# Pilot stats cards
stats_data = [
    ('🏫', '3个校区', '覆盖全部校区'),
    ('📍', '47个点位', '消防/电气/建筑/食堂'),
    ('👥', '6名巡检员', '保安人员（5人>45岁）'),
    ('🕐', '2个月', '2025年3月-4月'),
]

for i, (icon, value, label) in enumerate(stats_data):
    left = Inches(1.5 + i * 2.7)
    top = Inches(1.9)
    add_card(slide, left, top, Inches(2.4), Inches(2.2), NAVY_LIGHT)
    add_text_box(slide, left, top + Inches(0.3), Inches(2.4), Inches(0.6),
                 icon, font_size=40, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(1.0), Inches(2.4), Inches(0.5),
                 value, font_size=22, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(1.5), Inches(2.4), Inches(0.4),
                 label, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Training feedback
add_card(slide, Inches(1.5), Inches(4.5), Inches(5.3), Inches(2.5), NAVY_LIGHT)
add_text_box(slide, Inches(1.8), Inches(4.7), Inches(4.5), Inches(0.5),
             '📊 培训情况', font_size=18, color=CYAN, bold=True)
add_text_box(slide, Inches(1.8), Inches(5.3), Inches(4.8), Inches(0.5),
             '培训时长：约30分钟集中培训', font_size=14, color=WHITE_SOFT)
add_text_box(slide, Inches(1.8), Inches(5.8), Inches(4.8), Inches(0.8),
             '培训结果：6人全部可独立操作\n关键因素：微信小程序界面认知零门槛，操作流程简化为三步',
             font_size=13, color=GRAY)

add_card(slide, Inches(7.3), Inches(4.5), Inches(5.3), Inches(2.5), NAVY_LIGHT)
add_text_box(slide, Inches(7.6), Inches(4.7), Inches(4.5), Inches(0.5),
             '💡 用户体验反馈', font_size=18, color=CYAN, bold=True)
add_text_box(slide, Inches(7.6), Inches(5.3), Inches(4.8), Inches(0.5),
             '操作路径：选点位 → 拍照 → 提交（三步完成）', font_size=14, color=WHITE_SOFT)
add_text_box(slide, Inches(7.6), Inches(5.8), Inches(4.8), Inches(0.8),
             '巡检员评价："不用记哪个巡了哪个没巡"\n管理者评价："导出Excel对接检查省了大事"',
             font_size=13, color=GRAY)

# ═══════════════════════════════════════════════
# SLIDE 11 — Application Effects
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '应用效果 — 数据对比', font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(0.4),
             '2025年3-4月（数字化） vs 2024年10-11月（纸质）同期对比',
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Effect cards
effects = [
    ('📈', '巡检覆盖率', '纸质时期：记录丢失率~22%\n数字化后：系统强制留痕，基本全覆盖', GREEN_ACCENT),
    ('🔍', '隐患发现量', '纸质时期：月均十余条\n数字化后：月均三十余条\n原来被忽略的问题被纳入管理', CYAN_GLOW),
    ('🎯', '判定准确率', '人工判定误判率约35%\nAI辅助后下降至约8%\n有客观参考标准，减少分歧', GREEN_ACCENT),
    ('⚡', '整改进度', '平均周期：5.8天 → 1.6天\n每个环节有记录和责任人\n大幅减少沟通和推诿成本', CYAN_GLOW),
    ('⏱️', '管理效率', '每周汇总：3-4小时 → <1分钟\n一键导出Excel\n对接上级检查更加从容', GREEN_ACCENT),
]

for i, (icon, title, desc, accent) in enumerate(effects):
    row = i // 3
    col = i % 3
    if i >= 3:
        left = Inches(2.5 + col * 3.6)
        top = Inches(4.6)
    else:
        left = Inches(1.5 + col * 3.6)
        top = Inches(2.3)

    add_card(slide, left, top, Inches(3.2), Inches(2.0), NAVY_LIGHT)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(2.8), Inches(0.35),
                 f'{icon} {title}', font_size=16, color=accent, bold=True)
    tf = add_text_box(slide, left + Inches(0.2), top + Inches(0.6), Inches(2.8), Inches(1.2),
                      '', font_size=12, color=WHITE_SOFT)
    lines = desc.split('\n')
    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        add_rich_text(tf, line, font_size=12, color=GRAY)

# ═══════════════════════════════════════════════
# SLIDE 12 — Case Studies
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '应用效果 — 典型案例', font_size=32, color=WHITE, bold=True)

# Case 1
add_card(slide, Inches(1.5), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(1.8), Inches(2.1), Inches(4.5), Inches(0.5),
             '案例一：墙体裂缝 — AI发现结构隐患', font_size=18, color=RED_ACCENT, bold=True)

case1_steps = [
    '巡检员拍摄墙面裂缝照片',
    'AI判定：重大隐患（建筑安全）',
    '起初认为"不算严重"，AI坚持标记',
    '请建筑专业人士复查',
    '确认：确实为结构性裂缝',
    '→ 及时加固，避免了安全事故',
]
for i, step in enumerate(case1_steps):
    prefix = '🔴' if i == 5 else '→'
    add_text_box(slide, Inches(1.8), Inches(2.8 + i * 0.55), Inches(4.8), Inches(0.5),
                 f'{prefix} {step}', font_size=14,
                 color=WHITE_SOFT if i == 5 else GRAY)

# Case 2
add_card(slide, Inches(7.3), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(7.6), Inches(2.1), Inches(4.5), Inches(0.5),
             '案例二：消防通道 — AI发现被忽视的隐患', font_size=18, color=ORANGE_ACCENT, bold=True)

case2_steps = [
    '巡检员拍摄走廊堆放的纸箱',
    'AI判定：消防隐患（消防安全）',
    '巡检员最初不以为然：放点东西而已',
    '系统反复提醒后清理了通道',
    '事后消防演练证实该通道确实关键',
    '→ 帮助注意到了被主观忽略的细节',
]
for i, step in enumerate(case2_steps):
    prefix = '🟠' if i == 5 else '→'
    add_text_box(slide, Inches(7.6), Inches(2.8 + i * 0.55), Inches(4.8), Inches(0.5),
                 f'{prefix} {step}', font_size=14,
                 color=WHITE_SOFT if i == 5 else GRAY)

# Bottom
add_text_box(slide, Inches(2), Inches(7.0), Inches(9), Inches(0.4),
             'AI的价值：不是替代人做判断，而是提供一个不疲劳、不偷懒、有标准的"第二双眼睛"',
             font_size=14, color=CYAN_GLOW, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 13 — Innovation Points
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '创新点', font_size=32, color=WHITE, bold=True)

innovations = [
    ('01', '国产大模型 × 垂直场景', '针对校园安全巡检定义了6种隐患类型和2个等级标准\n结构化提示词 + AI结果自动回填 + 人工确认修正\n"AI辅助判定，人工最终把关"，兼顾效率与责任'),
    ('02', '完整业务闭环', '发现→AI判定→派单→整改→验收，全流程覆盖\n每个环节有时间戳、操作人和佐证照片\n可追溯、可审计，拒绝碎片化管理'),
    ('03', '低门槛低成本的实用方案', '微信小程序零安装，30分钟即可培训上岗\nWindows服务器部署，学校现有环境复用\nAI分析单次约0.002元，中小学预算完全可控'),
    ('04', 'AI辅助开发的示范意义', '非计算机专业教师借助AI完成全栈三端系统\n从数据库到API到前端界面，跨越技能限制\n"AI赋能终身学习"的具体实践案例'),
]

for i, (num, title, desc) in enumerate(innovations):
    top = Inches(1.8 + i * 1.35)
    # Number
    add_circle(slide, Inches(1.5), top + Inches(0.15), Inches(0.6), CYAN)
    add_text_box(slide, Inches(1.5), top + Inches(0.18), Inches(0.6), Inches(0.5),
                 num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Content
    add_card(slide, Inches(2.4), top, Inches(9.5), Inches(1.15), NAVY_LIGHT)
    add_text_box(slide, Inches(2.7), top + Inches(0.1), Inches(8.8), Inches(0.35),
                 title, font_size=20, color=CYAN, bold=True)
    tf = add_text_box(slide, Inches(2.7), top + Inches(0.45), Inches(8.8), Inches(0.65),
                      '', font_size=13, color=GRAY)
    lines = desc.split('\n')
    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        add_rich_text(tf, line, font_size=13, color=GRAY_LIGHT)

# ═══════════════════════════════════════════════
# SLIDE 14 — Problems & Improvements
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_cyan_bar(slide, Inches(0.8), Inches(0.8), Inches(0.4), Pt(4))

add_text_box(slide, Inches(1.5), Inches(0.7), Inches(10), Inches(0.7),
             '问题与改进方向', font_size=32, color=WHITE, bold=True)

# Problems (left)
add_card(slide, Inches(1.5), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(1.8), Inches(2.1), Inches(4.5), Inches(0.5),
             '⚠️ 遇到的问题', font_size=20, color=ORANGE_ACCENT, bold=True)

problems = [
    ('AI误判', '光线暗/照片模糊时准确率下降\n短期：拍照引导提示\n长期：积累数据微调模型'),
    ('微信限制', '订阅消息字段字符数卡得很严\nthing类最多20字符，文案需要反复精简'),
    ('Token并发', '多个请求同时401时的竞态处理\n第一版未考虑，经多轮调试才稳定'),
    ('照片质量', '部分巡检员拍摄敷衍\n后续需增加照片清晰度自动检测'),
]
for i, (title, desc) in enumerate(problems):
    add_text_box(slide, Inches(1.8), Inches(2.8 + i * 0.95), Inches(4.8), Inches(0.25),
                 f'• {title}', font_size=15, color=WHITE_SOFT, bold=True)
    add_text_box(slide, Inches(2.1), Inches(3.08 + i * 0.95), Inches(4.5), Inches(0.7),
                 desc, font_size=12, color=GRAY)

# Improvements (right)
add_card(slide, Inches(7.3), Inches(1.9), Inches(5.3), Inches(4.8), NAVY_LIGHT)
add_text_box(slide, Inches(7.6), Inches(2.1), Inches(4.5), Inches(0.5),
             '🚀 未来改进计划', font_size=20, color=CYAN, bold=True)

improvements = [
    ('短期', [
        'AI趋势预警：识别高频问题点位和季节风险',
        '语音输入：手套作业等不便打字场景',
        'GitHub开源：推动更多学校低成本使用',
    ]),
    ('远期', [
        '物联网联动：烟感/温感 + 人工巡检',
        '模型微调：积累标注数据训练专用模型',
        '"自动监测+人工巡检"双层安全体系',
    ]),
]

for period, items in improvements:
    add_text_box(slide, Inches(7.6), Inches(2.7 if period == '短期' else 4.3), Inches(4.5), Inches(0.3),
                 f'📅 {period}计划', font_size=16, color=CYAN if period == '短期' else CYAN_GLOW, bold=True)
    for j, item in enumerate(items):
        y = (3.1 if period == '短期' else 4.7) + j * 0.55
        add_text_box(slide, Inches(7.8), y, Inches(4.5), Inches(0.45),
                     f'→ {item}', font_size=13, color=WHITE_SOFT)

# ═══════════════════════════════════════════════
# SLIDE 15 — Thank You + AI Evidence
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY_DEEP)

# Decorative
add_circle(slide, Inches(-1), Inches(5), Inches(3), RGBColor(0x06, 0xB6, 0xD4))
add_circle(slide, Inches(10.5), Inches(-1), Inches(3.5), RGBColor(0x1A, 0x27, 0x44))
add_cyan_bar(slide, Inches(5.8), Inches(2.6), Inches(1.5), Pt(4))

# Center content
add_text_box(slide, Inches(2), Inches(1.5), Inches(9), Inches(0.8),
             '感谢评审', font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(2.8), Inches(9), Inches(0.6),
             '智巡校安 — 基于国产大模型的校园安全智能巡检信息系统', font_size=22, color=CYAN_GLOW,
             alignment=PP_ALIGN.CENTER)

# AI evidence summary
add_card(slide, Inches(2.5), Inches(3.8), Inches(8.3), Inches(2.5), NAVY_LIGHT)
add_text_box(slide, Inches(2.8), Inches(4.0), Inches(7.5), Inches(0.5),
             'AI辅助开发证据', font_size=18, color=CYAN, bold=True)

ai_items = [
    'Claude参与了全局CSS设计系统、AI提示词框架、微信推送流程、Token刷新机制等模块的生成',
    '所有AI生成代码均经人工审查、调试和多轮修改，关键模块（权限、校验、限速、部署）为人工编写',
    'AI生成代码已在源文件中以注释方式标注，便于评审查证',
]
for i, item in enumerate(ai_items):
    add_text_box(slide, Inches(2.8), Inches(4.6 + i * 0.55), Inches(7.5), Inches(0.5),
                 f'• {item}', font_size=13, color=GRAY)

# Bottom
add_text_box(slide, Inches(2), Inches(6.6), Inches(9), Inches(0.5),
             '配套资源：完整源代码 · 部署文档 · 演示视频 · 安装手册',
             font_size=14, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, '智巡校安_竞赛演示.pptx')
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')
